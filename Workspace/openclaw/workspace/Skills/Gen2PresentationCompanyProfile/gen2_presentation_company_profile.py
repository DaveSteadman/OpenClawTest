#!/usr/bin/env python3
"""Gen2PresentationCompanyProfile - Build a CompanyProfile slide from analysis JSON."""

import argparse
import datetime as dt
import json
import os
import re
import sys
from pathlib import Path

SKILLS_ROOT = Path(__file__).resolve().parent.parent
if str(SKILLS_ROOT) not in sys.path:
    sys.path.insert(0, str(SKILLS_ROOT))

from CommonCode.FolderNavigator import FolderNavigator
from CommonCode.PresentationMakerCompanyProfile import instantiate_template

try:
    import requests
    from pptx import Presentation
except ImportError:
    print(
        json.dumps(
            {
                "domain": "",
                "date": "",
                "error": "Missing dependencies. Install with: pip install -r requirements.txt",
                "status": "error",
            },
            ensure_ascii=False,
        )
    )
    sys.exit(1)


def validate_domain(domain: str) -> str:
    if not domain or not isinstance(domain, str):
        raise ValueError("Domain must be a non-empty string")
    cleaned = domain.strip()
    if not re.fullmatch(r"[A-Za-z]+", cleaned):
        raise ValueError("Domain must be alphabetic only (A-Z, a-z)")
    return cleaned


def parse_date(value: str) -> dt.date:
    cleaned = (value or "").strip().replace("-", "/")
    parts = cleaned.split("/")
    if len(parts) != 3:
        raise ValueError("Date must be YYYY/MM/DD or YYYY-MM-DD")
    try:
        year, month, day = (int(item) for item in parts)
        return dt.date(year, month, day)
    except ValueError as exc:
        raise ValueError("Invalid date value") from exc


def safe_name(text: str, fallback: str = "CompanyProfile") -> str:
    tokens = re.findall(r"[A-Za-z0-9]+", text or "")
    if not tokens:
        return fallback
    parts = []
    for token in tokens:
        if token.isdigit():
            parts.append(token)
        else:
            parts.append(token[0].upper() + token[1:].lower())
    return "".join(parts)[:80]


def find_analysis_files(navigator: FolderNavigator, domain: str, as_of_date: dt.date):
    domain_root = navigator.get_domain_root("analyse", domain)
    if not domain_root.exists():
        raise FileNotFoundError(f"Analysis domain folder not found: {domain_root}")

    day_dir = domain_root / f"{as_of_date.year:04d}" / f"{as_of_date.month:02d}" / f"{as_of_date.day:02d}"
    month_dir = domain_root / f"{as_of_date.year:04d}" / f"{as_of_date.month:02d}"
    year_dir = domain_root / f"{as_of_date.year:04d}"

    for folder in (day_dir, month_dir, year_dir):
        if folder.exists() and folder.is_dir():
            files = sorted(folder.glob("*.json"))
            if files:
                return files

    raise FileNotFoundError(
        f"No analysis JSON files found for domain '{domain}' near date '{as_of_date.isoformat()}'"
    )


def _stringify_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, list):
        return " | ".join(str(item) for item in value if str(item).strip())
    return str(value).strip()


def build_context_from_analysis(files: list[Path], max_files: int, max_chars: int):
    selected = files[: max(1, max_files)]
    blocks = []
    total = 0
    for path in selected:
        try:
            payload = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        except Exception:
            continue

        analysis = payload.get("analysis", payload)
        if not isinstance(analysis, dict):
            analysis = {"raw": analysis}

        summary = _stringify_value(analysis.get("executive_summary"))
        trends = _stringify_value(analysis.get("key_trends"))
        delta = _stringify_value(analysis.get("delta_start_to_end"))
        evidence = _stringify_value(analysis.get("supporting_evidence"))
        caveats = _stringify_value(analysis.get("caveats"))
        confidence = _stringify_value(analysis.get("confidence"))

        block = (
            f"### Source: {path.name}\n"
            f"Summary: {summary}\n"
            f"Trends: {trends}\n"
            f"Delta: {delta}\n"
            f"Evidence: {evidence}\n"
            f"Caveats: {caveats}\n"
            f"Confidence: {confidence}"
        ).strip()

        next_total = total + len(block)
        if max_chars > 0 and next_total > max_chars:
            break

        blocks.append(block)
        total = next_total

    if not blocks:
        raise RuntimeError("No readable analysis content after bounds were applied")

    return selected, "\n\n".join(blocks)


def call_llm_slide_formatter(
    domain: str,
    as_of_date: dt.date,
    prompt_text: str,
    analysis_context: str,
    llm_timeout: int,
):
    api_key = os.environ.get("OPENAI_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is required for LLM formatting")

    model = os.environ.get("OPENAI_MODEL", "gpt-4o-mini").strip() or "gpt-4o-mini"
    base_url = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1").strip().rstrip("/")
    endpoint = f"{base_url}/chat/completions"

    system_prompt = (
        "You are a presentation editor. Convert analysis into concise executive slide content. "
        "Return strict JSON only with no markdown fences and no extra keys."
    )

    user_prompt = (
        f"Domain: {domain}\n"
        f"Date: {as_of_date.isoformat()}\n"
        f"Instruction: {prompt_text}\n\n"
        "Return strict JSON with keys:\n"
        "- slide_title (string)\n"
        "- main_points (array of 4 to 10 strings, concise and factual)\n"
        "- source_note (string)\n\n"
        "Rules:\n"
        "- Prioritize money-linked events when available.\n"
        "- Do not invent numbers; if uncertain, use wording like 'reported'.\n"
        "- Keep each point under 120 characters when possible.\n\n"
        "Analysis context:\n"
        f"{analysis_context}"
    )

    payload = {
        "model": model,
        "temperature": 0.2,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    }

    response = requests.post(
        endpoint,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        json=payload,
        timeout=max(30, llm_timeout),
    )
    response.raise_for_status()

    body = response.json()
    content = body.get("choices", [{}])[0].get("message", {}).get("content", "")
    if not content:
        raise RuntimeError("LLM returned empty response")

    cleaned = content.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?", "", cleaned).strip()
        cleaned = re.sub(r"```$", "", cleaned).strip()

    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise RuntimeError("LLM response was not valid JSON") from exc

    return {
        "model": model,
        "endpoint": endpoint,
        "payload": parsed,
        "raw_response": cleaned,
    }


def normalize_slide_payload(data: dict, domain: str, as_of_date: dt.date):
    title = str(data.get("slide_title") or f"{domain} Company Profile").strip()
    source_note = str(data.get("source_note") or "Generated from analysis artifacts").strip()

    points_raw = data.get("main_points")
    points = []
    if isinstance(points_raw, list):
        for item in points_raw:
            text = str(item).strip()
            if not text:
                continue
            if len(text) > 130:
                text = text[:127].rstrip() + "..."
            points.append(text)

    if not points:
        points = [
            "No clear major events extracted from available analysis.",
            "Review source analysis files and rerun with expanded date coverage.",
        ]

    if len(points) > 10:
        points = points[:10]

    main_text = "\n".join(f"• {item}" for item in points)

    return {
        "slide_title": title,
        "main_text": main_text,
        "source_note": source_note,
        "date": as_of_date.isoformat(),
    }


def _set_shape_text(shape, text: str):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.text = text or ""


def _find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def render_company_profile_slide_from_template(
    out_dir: Path,
    domain: str,
    as_of_date: dt.date,
    slide_data: dict,
    source_files: list[Path],
):
    out_dir.mkdir(parents=True, exist_ok=True)

    base_name = safe_name(slide_data["slide_title"], fallback=f"{domain}CompanyProfile")
    scope = as_of_date.strftime("%Y-%m-%d")
    pptx_path = out_dir / f"{base_name}_{scope}.pptx"

    instantiate_template(output_dir=out_dir, output_name=pptx_path.name, overwrite=True)

    prs = Presentation(str(pptx_path))
    if not prs.slides:
        raise RuntimeError("Template did not contain a slide")

    slide = prs.slides[0]
    title_shape = _find_shape_by_name(slide, "FIELD_TITLE")
    main_text_shape = _find_shape_by_name(slide, "FIELD_MAIN_TEXT")
    date_shape = _find_shape_by_name(slide, "FIELD_DATE")
    source_shape = _find_shape_by_name(slide, "FIELD_SOURCE_NOTE")

    if title_shape is None or main_text_shape is None:
        raise RuntimeError("Template missing required named fields: FIELD_TITLE and FIELD_MAIN_TEXT")

    _set_shape_text(title_shape, slide_data["slide_title"])
    _set_shape_text(main_text_shape, slide_data["main_text"])
    if date_shape is not None:
        _set_shape_text(date_shape, slide_data["date"])
    if source_shape is not None:
        _set_shape_text(source_shape, f"Sources: {len(source_files)} | {slide_data['source_note']}")

    prs.save(str(pptx_path))

    payload_path = out_dir / f"{base_name}_{scope}.json"
    payload = {
        "ts": dt.datetime.now(dt.timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z"),
        "domain": domain,
        "date": as_of_date.isoformat(),
        "source_files": [str(path) for path in source_files],
        "slide": slide_data,
    }
    payload_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    return pptx_path, payload_path


def main():
    parser = argparse.ArgumentParser(description="Generate CompanyProfile PPTX slide from analysis JSON")
    parser.add_argument("domain", help="Domain name (alphabetic only)")
    parser.add_argument("--date", required=True, help="Date in YYYY/MM/DD or YYYY-MM-DD")
    parser.add_argument(
        "--prompt",
        default="Format this analysis into CompanyProfile slide bullets.",
        help="LLM instruction for bullet formatting",
    )
    parser.add_argument("--analysis-file", default="", help="Optional path to a single analysis JSON")
    parser.add_argument("--max-analysis-files", type=int, default=6)
    parser.add_argument("--max-analysis-chars", type=int, default=18000)
    parser.add_argument("--llm-timeout", type=int, default=90)

    try:
        args = parser.parse_args()
    except SystemExit:
        return 1

    try:
        domain = validate_domain(args.domain)
        as_of_date = parse_date(args.date)
        prompt_text = (args.prompt or "").strip()
        if not prompt_text:
            raise ValueError("Prompt must be non-empty")
    except Exception as exc:
        print(
            json.dumps(
                {
                    "domain": args.domain if hasattr(args, "domain") else "",
                    "date": args.date if hasattr(args, "date") else "",
                    "error": str(exc),
                    "status": "error",
                },
                ensure_ascii=False,
            )
        )
        return 1

    navigator = FolderNavigator.from_fixed_point()

    try:
        if args.analysis_file:
            candidate = Path(args.analysis_file)
            if not candidate.is_absolute():
                candidate = Path.cwd() / candidate
            if not candidate.exists():
                raise FileNotFoundError(f"Analysis file not found: {candidate}")
            files = [candidate]
        else:
            files = find_analysis_files(navigator=navigator, domain=domain, as_of_date=as_of_date)

        selected_files, context = build_context_from_analysis(
            files=files,
            max_files=max(1, args.max_analysis_files),
            max_chars=max(1000, args.max_analysis_chars),
        )

        llm_result = call_llm_slide_formatter(
            domain=domain,
            as_of_date=as_of_date,
            prompt_text=prompt_text,
            analysis_context=context,
            llm_timeout=max(30, args.llm_timeout),
        )

        slide_data = normalize_slide_payload(
            data=llm_result["payload"],
            domain=domain,
            as_of_date=as_of_date,
        )

        out_dir = navigator.get_date_path(
            area="present",
            domain=domain,
            value=as_of_date,
            create=True,
        )

        presentation_path, payload_path = render_company_profile_slide_from_template(
            out_dir=out_dir,
            domain=domain,
            as_of_date=as_of_date,
            slide_data=slide_data,
            source_files=selected_files,
        )

        print(
            json.dumps(
                {
                    "domain": domain,
                    "date": as_of_date.isoformat(),
                    "analysis_files": [str(path) for path in selected_files],
                    "presentation_path": str(presentation_path),
                    "payload_path": str(payload_path),
                    "status": "ok",
                },
                ensure_ascii=False,
            )
        )
        return 0

    except Exception as exc:
        print(
            json.dumps(
                {
                    "domain": domain,
                    "date": as_of_date.isoformat(),
                    "error": str(exc),
                    "status": "error",
                },
                ensure_ascii=False,
            )
        )
        return 1


if __name__ == "__main__":
    sys.exit(main())
