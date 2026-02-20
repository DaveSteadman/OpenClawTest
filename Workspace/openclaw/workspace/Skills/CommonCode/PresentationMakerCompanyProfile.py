#!/usr/bin/env python3
"""Create a CompanyProfile PowerPoint template with named fields."""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    print(
        json.dumps(
            {
                "error": "Missing dependency python-pptx. Install with: pip install python-pptx",
                "status": "error",
            },
            ensure_ascii=False,
        )
    )
    sys.exit(1)


FIELD_TOKENS = {
    "FIELD_TITLE": "{{TITLE}}",
    "FIELD_LOGO": "{{LOGO}}",
    "FIELD_DATE": "{{DATE}}",
    "FIELD_MAIN_TEXT": "{{MAIN_TEXT}}",
    "FIELD_SOURCE_NOTE": "{{SOURCE_NOTE}}",
}

TEMPLATE_FILENAME = "CompanyProfile.pptx"


def _set_shape_name(shape, name: str) -> None:
    shape.element.nvSpPr.cNvPr.set("name", name)


def create_company_profile_template(target_path: str | Path) -> Path:
    out_path = Path(target_path).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titlebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(11.69), Inches(1.1))
    _set_shape_name(titlebar, "DECOR_TITLEBAR")
    titlebar.line.fill.background()
    titlebar.fill.solid()
    titlebar.fill.fore_color.rgb = RGBColor(22, 44, 84)

    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.0), Inches(0.6))
    _set_shape_name(title_shape, "FIELD_TITLE")
    title_tf = title_shape.text_frame
    title_tf.clear()
    title_run = title_tf.paragraphs[0]
    title_run.text = FIELD_TOKENS["FIELD_TITLE"]
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)

    logo_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.25), Inches(0.2), Inches(2.1), Inches(0.75))
    _set_shape_name(logo_shape, "FIELD_LOGO")
    logo_shape.fill.background()
    logo_shape.line.color.rgb = RGBColor(255, 255, 255)
    logo_tf = logo_shape.text_frame
    logo_tf.clear()
    logo_run = logo_tf.paragraphs[0]
    logo_run.text = FIELD_TOKENS["FIELD_LOGO"]
    logo_run.font.size = Pt(12)
    logo_run.font.bold = True
    logo_run.font.color.rgb = RGBColor(255, 255, 255)

    date_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2.4), Inches(0.3))
    _set_shape_name(date_shape, "FIELD_DATE")
    date_tf = date_shape.text_frame
    date_tf.clear()
    date_run = date_tf.paragraphs[0]
    date_run.text = FIELD_TOKENS["FIELD_DATE"]
    date_run.font.size = Pt(11)

    main_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.55), Inches(10.8), Inches(5.85))
    _set_shape_name(main_shape, "FIELD_MAIN_TEXT")
    main_shape.fill.solid()
    main_shape.fill.fore_color.rgb = RGBColor(246, 248, 252)
    main_shape.line.color.rgb = RGBColor(195, 203, 217)

    main_tf = main_shape.text_frame
    main_tf.word_wrap = True
    main_tf.clear()
    main_run = main_tf.paragraphs[0]
    main_run.text = FIELD_TOKENS["FIELD_MAIN_TEXT"]
    main_run.font.size = Pt(20)

    source_shape = slide.shapes.add_textbox(Inches(0.5), Inches(7.65), Inches(10.8), Inches(0.3))
    _set_shape_name(source_shape, "FIELD_SOURCE_NOTE")
    source_tf = source_shape.text_frame
    source_tf.clear()
    source_run = source_tf.paragraphs[0]
    source_run.text = FIELD_TOKENS["FIELD_SOURCE_NOTE"]
    source_run.font.size = Pt(10)

    prs.save(str(out_path))
    return out_path


def get_standard_template_path() -> Path:
    return Path(__file__).resolve().parent / TEMPLATE_FILENAME


def ensure_standard_template(force: bool = False) -> Path:
    template_path = get_standard_template_path()
    if force or not template_path.exists():
        create_company_profile_template(template_path)
    return template_path


def instantiate_template(output_dir: str | Path, output_name: str = TEMPLATE_FILENAME, overwrite: bool = True) -> Path:
    template_path = ensure_standard_template(force=False)
    target_dir = Path(output_dir).resolve()
    target_dir.mkdir(parents=True, exist_ok=True)
    target_path = target_dir / (output_name or TEMPLATE_FILENAME)
    if target_path.exists() and not overwrite:
        raise FileExistsError(f"Target file already exists: {target_path}")
    shutil.copyfile(template_path, target_path)
    return target_path


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Ensure standard CompanyProfile template and optionally instantiate a copy"
    )
    parser.add_argument(
        "--force-template",
        action="store_true",
        help="Rebuild the standard template next to this Python file",
    )
    parser.add_argument(
        "--instantiate-dir",
        default="",
        help="Optional directory to copy a fresh template instance into",
    )
    parser.add_argument(
        "--output-name",
        default=TEMPLATE_FILENAME,
        help="Output filename when using --instantiate-dir",
    )
    parser.add_argument(
        "--no-overwrite",
        action="store_true",
        help="Do not overwrite existing instantiated output file",
    )

    try:
        args = parser.parse_args()
        template_path = ensure_standard_template(force=bool(args.force_template))
        instantiated_path = ""
        if args.instantiate_dir:
            instantiated = instantiate_template(
                output_dir=args.instantiate_dir,
                output_name=args.output_name,
                overwrite=not bool(args.no_overwrite),
            )
            instantiated_path = str(instantiated)

        print(
            json.dumps(
                {
                    "template_path": str(template_path),
                    "instantiated_path": instantiated_path,
                    "fields": FIELD_TOKENS,
                    "status": "ok",
                },
                ensure_ascii=False,
            )
        )
        return 0
    except Exception as exc:
        print(json.dumps({"error": str(exc), "status": "error"}, ensure_ascii=False))
        return 1


if __name__ == "__main__":
    sys.exit(main())
